import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set 16:9 Widescreen slides
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette Constants
    BG_COLOR = RGBColor(11, 12, 16)      # #0b0c10 Deep Black
    PANEL_BG = RGBColor(26, 29, 38)      # #1a1d26 Dark Slate Panel
    BORDER_COLOR = RGBColor(40, 44, 55)  # #282c37 Muted Border
    ACCENT_CYAN = RGBColor(0, 240, 255)  # #00f0ff Electric Cyan
    ACCENT_ORANGE = RGBColor(255, 107, 0) # #ff6b00 Vibrant Orange
    TEXT_MAIN = RGBColor(245, 246, 250)  # #f5f6fa Off-White
    TEXT_MUTED = RGBColor(154, 160, 176) # #9aa0b0 Grey Muted

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    def set_dark_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_title(slide, category, title, desc=""):
        # Add textbox for category and title
        tb = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        # Category
        p1 = tf.paragraphs[0]
        p1.text = category.upper()
        p1.font.size = Pt(10)
        p1.font.color.rgb = ACCENT_CYAN
        p1.font.bold = True
        p1.font.name = 'Space Grotesk'
        
        # Title
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_MAIN
        p2.font.name = 'Space Grotesk'
        p2.space_before = Pt(4)
        
        if desc:
            p3 = tf.add_paragraph()
            p3.text = desc
            p3.font.size = Pt(13)
            p3.font.color.rgb = TEXT_MUTED
            p3.font.name = 'Arial'
            p3.space_before = Pt(4)

    def add_card(slide, left, top, width, height, title, value, unit="", color_accent=ACCENT_CYAN):
        # Draw background panel
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = PANEL_BG
        shape.line.color.rgb = BORDER_COLOR
        shape.line.width = Pt(1)
        
        # Text
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.18)
        tf.margin_top = Inches(0.18)
        
        p = tf.paragraphs[0]
        p.text = title.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = TEXT_MUTED
        p.font.name = 'Arial'
        
        p2 = tf.add_paragraph()
        p2.text = f"{value} {unit}".strip()
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = color_accent
        p2.font.name = 'Space Grotesk'
        p2.space_before = Pt(6)

    # ==========================================
    # SLIDE 1: PORTADA
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_dark_background(slide1)

    # Title box (Left Side)
    tb1 = slide1.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(6.5), Inches(4.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_top = 0

    p_cat = tf1.paragraphs[0]
    p_cat.text = "PROYECTO ACADÉMICO - TURBOMÁQUINAS"
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_CYAN
    p_cat.font.name = 'Space Grotesk'

    p_title = tf1.add_paragraph()
    p_title.text = "Diseño de Ventilador\nCentrífugo Sirocco"
    p_title.font.size = Pt(40)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_MAIN
    p_title.font.name = 'Space Grotesk'
    p_title.space_before = Pt(14)
    p_title.space_after = Pt(14)

    p_desc = tf1.add_paragraph()
    p_desc.text = "Dimensionamiento reproducible, triangulación de velocidades y modelado técnico de un ventilador multialabe de doble ancho para aire limpio."
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = TEXT_MUTED
    p_desc.font.name = 'Arial'
    
    p_group = tf1.add_paragraph()
    p_group.text = "Desarrollado por: Grupo 2\nRevisión: Versión V1 (Cierre Unidimensional)"
    p_group.font.size = Pt(12)
    p_group.font.bold = True
    p_group.font.color.rgb = ACCENT_ORANGE
    p_group.font.name = 'Space Grotesk'
    p_group.space_before = Pt(28)

    # Image (Right Side)
    img_path = "assets/sirocco_fan_render.jpg"
    if os.path.exists(img_path):
        slide1.shapes.add_picture(img_path, Inches(7.5), Inches(1.5), width=Inches(5.0), height=Inches(4.5))
    else:
        # Placeholder shape if image is missing
        shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(1.5), Inches(5.0), Inches(4.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PANEL_BG
        shape.line.color.rgb = BORDER_COLOR
        tx = shape.text_frame
        tx.text = "[Render del Ventilador]"
        tx.paragraphs[0].font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 2: REQUERIMIENTOS
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_dark_background(slide2)
    add_title(slide2, "Requerimientos de Diseño", "Parámetros de Entrada y Configuración", "Punto de operación establecido para el dimensionamiento del sistema.")

    # 3 Metrics Cards
    add_card(slide2, Inches(0.75), Inches(2.0), Inches(3.6), Inches(1.5), "Caudal Solicitado (Q)", "3.0", "m³/s")
    add_card(slide2, Inches(4.6), Inches(2.0), Inches(3.6), Inches(1.5), "Presión Estática (Δp_s)", "90", "mmH₂O", ACCENT_CYAN)
    add_card(slide2, Inches(8.45), Inches(2.0), Inches(4.1), Inches(1.5), "Densidad de Diseño (ρ)", "1.20", "kg/m³")

    # Bottom Panel: DIDW Justification
    shape2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(3.9), Inches(11.8), Inches(2.8))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = PANEL_BG
    shape2.line.color.rgb = BORDER_COLOR
    
    tf2 = shape2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = Inches(0.25)
    
    p = tf2.paragraphs[0]
    p.text = "JUSTIFICACIÓN DE LA SELECCIÓN DE DOBLE ENTRADA (DIDW)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.font.name = 'Space Grotesk'
    
    p2 = tf2.add_paragraph()
    p2.text = "• Un caudal de 3.0 m³/s (10,800 m³/h) es elevado para un ventilador de entrada simple (SISW), lo que obligaría a diseñar un rodete con un ojo de succión sobredimensionado o a tolerar velocidades axiales muy altas.\n" \
              "• Al seleccionar un rodete de doble ancho y doble entrada (DIDW), el flujo de entrada se divide equitativamente (1.5 m³/s por cada succión), reduciendo la velocidad axial en el ojo a un valor óptimo de ~7.77 m/s.\n" \
              "• Esto minimiza las pérdidas por choque aerodinámico en la entrada de los álabes, mejora la eficiencia estática del rodete y disminuye significativamente la emisión de ruido tonal."
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_MUTED
    p2.font.name = 'Arial'
    p2.space_before = Pt(8)

    # ==========================================
    # SLIDE 3: GEOMETRÍA DEL ROTOR
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_dark_background(slide3)
    add_title(slide3, "Geometría del Rotor", "Especificación del Rodete Sirocco V1", "Parámetros geométricos para modelado 3D y planos de fabricación.")

    # Left Column: Specs Grid
    left_x = Inches(0.75)
    top_y = Inches(2.0)
    card_w = Inches(2.7)
    card_h = Inches(1.1)

    # Row 1
    add_card(slide3, left_x, top_y, card_w, card_h, "Diámetro Ext (D₂)", "600", "mm")
    add_card(slide3, left_x + Inches(2.9), top_y, card_w, card_h, "Diámetro Int (D₁)", "510", "mm")
    # Row 2
    add_card(slide3, left_x, top_y + Inches(1.3), card_w, card_h, "Ancho Rodete (b)", "230", "mm")
    add_card(slide3, left_x + Inches(2.9), top_y + Inches(1.3), card_w, card_h, "Número Álabes (Z)", "48", "álabes")
    # Row 3
    add_card(slide3, left_x, top_y + Inches(2.6), card_w, card_h, "Ángulo Salida (β₂)", "125", "°")
    add_card(slide3, left_x + Inches(2.9), top_y + Inches(2.6), card_w, card_h, "Ángulo Entrada (β₁)", "18", "°")
    # Row 4
    add_card(slide3, left_x, top_y + Inches(3.9), card_w, card_h, "Espesor Alabe (t)", "1.2", "mm")
    add_card(slide3, left_x + Inches(2.9), top_y + Inches(3.9), card_w, card_h, "Relación D₁/D₂", "0.85", "")

    # Right Column: Tech Panel
    shape3 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.0), Inches(5.75), Inches(5.0))
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = PANEL_BG
    shape3.line.color.rgb = BORDER_COLOR
    
    tf3 = shape3.text_frame
    tf3.word_wrap = True
    tf3.margin_left = tf3.margin_top = tf3.margin_right = tf3.margin_bottom = Inches(0.25)

    p = tf3.paragraphs[0]
    p.text = "CARACTERÍSTICAS DEL RODETE SIROCCO"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.font.name = 'Space Grotesk'
    
    p2 = tf3.add_paragraph()
    p2.text = "• Álabes Curvados hacia Adelante: El ángulo de salida beta2 = 125° (mayor a 90°) es característico de las turbomáquinas de acción que priorizan altas presiones en dimensiones compactas.\n\n" \
              "• Corona Circular Estrecha: La relación D1/D2 = 0.85 deja álabes radialmente muy cortos. Esto permite que el rodete actúe principalmente transfiriendo energía cinética (gran velocidad periférica U2).\n\n" \
              "• Espesor Fino de Chapa: El espesor preliminar de 1.2 mm requiere un alabeado por conformado de chapa. El bloqueo de flujo kb en la entrada por espesor se calculó en 12% (kb1 = 0.88), un valor aceptable.\n\n" \
              "• Transmisión y Motorización: Se estima una potencia mecánica al eje de 4.86 kW. Se adopta un motor comercial de 5.5 kW con un margen del 13.1%, conectado por correas con poleas de 125/200 mm para una relación de transmisión de ~0.63."
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_MUTED
    p2.font.name = 'Arial'
    p2.space_before = Pt(10)

    # ==========================================
    # SLIDE 4: TRIÁNGULO DE VELOCIDADES
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_dark_background(slide4)
    add_title(slide4, "Ingeniería de Velocidades", "Triángulo de Velocidades a 1100 RPM", "Cálculo y componentes de velocidad a la salida del rodete.")

    # Left Column: Metrics
    add_card(slide4, Inches(0.75), Inches(2.0), Inches(2.7), Inches(1.3), "Vel. Tangencial (U₂)", "34.56", "m/s")
    add_card(slide4, Inches(0.75), Inches(3.5), Inches(2.7), Inches(1.3), "Vel. Radial (C_m₂)", "7.19", "m/s")
    add_card(slide4, Inches(0.75), Inches(5.0), Inches(2.7), Inches(1.3), "Vel. Arrastre (C_u₂)", "37.51", "m/s")

    # Right Text Details Panel
    shape4 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.8), Inches(2.0), Inches(8.75), Inches(4.3))
    shape4.fill.solid()
    shape4.fill.fore_color.rgb = PANEL_BG
    shape4.line.color.rgb = BORDER_COLOR
    
    tf4 = shape4.text_frame
    tf4.word_wrap = True
    tf4.margin_left = tf4.margin_top = tf4.margin_right = tf4.margin_bottom = Inches(0.25)

    p = tf4.paragraphs[0]
    p.text = "ANÁLISIS DE LOS VECTORES DE VELOCIDAD"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.font.name = 'Space Grotesk'
    
    p2 = tf4.add_paragraph()
    p2.text = "• Deslizamiento de Wiesner: Debido al número finito de álabes (Z=48), el flujo no sigue perfectamente el ángulo constructivo. Se aplicó el factor de deslizamiento de Wiesner (sigma = 0.940), resultando en una velocidad tangencial real de salida Cu2 = 37.51 m/s.\n\n" \
              "• Velocidad Absoluta (C₂): Con un vector tangencial Cu2 de 37.51 m/s y una componente radial Cm2 de 7.19 m/s, la velocidad absoluta de salida del aire del rodete es C2 = 38.2 m/s, saliendo con una dirección inclinada muy pronunciada.\n\n" \
              "• Presión Teórica de Euler (Δp_E): El trabajo impartido teóricamente por el rodete al aire es de 1555.5 Pa (158.6 mmH2O).\n\n" \
              "• Eficiencia Estática (η_estática): Considerando las pérdidas internas y una eficiencia mecánica del 96% en la transmisión, la potencia al eje de cálculo es de 4.86 kW, lo que arroja una eficiencia estática estimada del 54.5% para el ventilador en este punto de diseño."
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = TEXT_MUTED
    p2.font.name = 'Arial'
    p2.space_before = Pt(10)

    # ==========================================
    # SLIDE 5: VOLUTA
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_dark_background(slide5)
    add_title(slide5, "Estructura del Casing", "Voluta y Conversión Dinámica", "Carcasa espiral diseñada para la transformación de velocidad en presión estática.")

    add_card(slide5, Inches(0.75), Inches(2.0), Inches(3.6), Inches(1.5), "Ancho Interior Voluta", "280", "mm")
    add_card(slide5, Inches(4.6), Inches(2.0), Inches(3.6), Inches(1.5), "Holgura de Lengua", "48", "mm", ACCENT_ORANGE)
    add_card(slide5, Inches(8.45), Inches(2.0), Inches(4.1), Inches(1.5), "Velocidad de Descarga", "~15", "m/s")

    # Bottom Panel
    shape5 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(3.9), Inches(11.8), Inches(2.8))
    shape5.fill.solid()
    shape5.fill.fore_color.rgb = PANEL_BG
    shape5.line.color.rgb = BORDER_COLOR
    
    tf5 = shape5.text_frame
    tf5.word_wrap = True
    tf5.margin_left = tf5.margin_top = tf5.margin_right = tf5.margin_bottom = Inches(0.25)

    p = tf5.paragraphs[0]
    p.text = "PRINCIPIOS DE DISEÑO DE LA VOLUTA"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.font.name = 'Space Grotesk'
    
    p2 = tf5.add_paragraph()
    p2.text = "• Conversión de Energía: El aire sale del rodete a 38.2 m/s (alta energía cinética). La voluta, diseñada siguiendo una espiral logarítmica, ensancha gradualmente el área de paso del flujo para disminuir la velocidad y aumentar la presión estática hasta alcanzar ~15 m/s en la boca de descarga.\n" \
              "• Holgura de la Lengua (Cut-off): Ubicada a 48 mm del rodete (8% de D2). Esta distancia evita la interacción de presiones severas que originarían ruido tonal (frecuencia de paso de álabes). Una holgura menor incrementaría ligeramente la presión, pero el ruido sería inaceptable para entornos industriales normales.\n" \
              "• Ancho de Carcasa: La voluta tiene 280 mm de ancho interior, lo que otorga una holgura lateral de 25 mm por lado respecto al rodete (230 mm) para evitar fricciones secundarias de disco y permitir el flujo de retorno."
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = TEXT_MUTED
    p2.font.name = 'Arial'
    p2.space_before = Pt(8)

    # ==========================================
    # SLIDE 6: COMPARATIVA
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_dark_background(slide6)
    add_title(slide6, "Análisis de Operación", "Comparativa de Velocidad de Giro", "Justificación del descarte de 1500 RPM frente a las 1100 RPM recomendadas.")

    # Table creation
    rows = 6
    cols = 4
    left = Inches(0.75)
    top = Inches(2.0)
    width = Inches(11.83)
    height = Inches(4.5)
    
    table_shape = slide6.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Set Column Widths
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(3.0)
    table.columns[2].width = Inches(3.0)
    table.columns[3].width = Inches(2.63)

    # Headers
    headers = ["MAGNITUD / PARÁMETRO", "OPERACIÓN A 1100 RPM (V1)", "OPERACIÓN A 1500 RPM (TANTEO)", "CONCLUSIÓN / ANÁLISIS"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL_BG
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN
        p.font.name = 'Space Grotesk'
        p.alignment = PP_ALIGN.CENTER
        
    data = [
        ["Velocidad Periférica (U₂)", "34.56 m/s", "47.12 m/s", "Mayor esfuerzo centrífugo a 1500 rpm"],
        ["Coeficiente de Presión", "0.616", "0.331", "Desequilibrio en coeficientes"],
        ["Presión de Euler (Δp_E)", "1555.5 Pa", "2892.4 Pa", "Exceso innecesario de energía"],
        ["Potencia al Eje Estimada", "4.86 kW", "9.03 kW", "Excede capacidad del motor de 5.5 kW"],
        ["Estado de Viabilidad", "VIABLE (Motor de 5.5 kW)", "INVIABLE (> 9 kW requeridos)", "Selección de 1100 rpm óptima"]
    ]
    
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL_BG if row_idx % 2 == 0 else BG_COLOR
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(10.5)
            p.font.name = 'Arial'
            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
            
            # Highlight highlights
            if "VIABLE" in text:
                p.font.color.rgb = RGBColor(0, 255, 102) # Green
                p.font.bold = True
            elif "INVIABLE" in text:
                p.font.color.rgb = ACCENT_ORANGE
                p.font.bold = True
            elif col_idx == 3:
                p.font.color.rgb = ACCENT_CYAN
            else:
                p.font.color.rgb = TEXT_MAIN

    # ==========================================
    # SLIDE 7: REPOSITORIO
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_dark_background(slide7)
    add_title(slide7, "Estructura del Proyecto", "Repaso del Repositorio de Código en GitHub", "Organización y contenidos de los archivos reproducibles del proyecto.")

    # Text list of repo directories
    shape7 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(2.0), Inches(11.83), Inches(4.5))
    shape7.fill.solid()
    shape7.fill.fore_color.rgb = PANEL_BG
    shape7.line.color.rgb = BORDER_COLOR
    
    tf7 = shape7.text_frame
    tf7.word_wrap = True
    tf7.margin_left = tf7.margin_top = tf7.margin_right = tf7.margin_bottom = Inches(0.25)

    p = tf7.paragraphs[0]
    p.text = "ENLACE: https://github.com/Aquiles-ryp360/diseno-ventilador-sirocco"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.font.name = 'Space Grotesk'
    
    dirs = [
        "📂 00_gestion - Cronogramas y actas de reuniones de diseño.",
        "📂 04_bibliografia - Catálogos de fabricantes y libros teóricos sobre ventiladores centrífugos.",
        "📂 05_avances - Informes académicos y reportes de progreso del diseño.",
        "📂 06_calculos - Memorias de cálculo (memoria_calculo_sirocco_v1.md) y verificaciones del eje.",
        "📂 07_planos - Archivos STEP CAD 3D de rotor y DXF 2D para el corte de voluta.",
        "📂 08_software - Scripts de simulación en Python / Octave para cálculo de espiral y velocidades.",
        "📂 presentacion - Diapositivas web interactivas (HTML, CSS, JS, imágenes)."
    ]
    
    for d in dirs:
        p_dir = tf7.add_paragraph()
        p_dir.text = d
        p_dir.font.size = Pt(12)
        p_dir.font.color.rgb = TEXT_MAIN
        p_dir.font.name = 'Arial'
        p_dir.space_before = Pt(8)

    # Save presentation
    output_filename = "presentacion_sirocco.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as {output_filename}")

if __name__ == "__main__":
    create_presentation()
